import aiohttp
import asyncio
import os
import whisper
import pypdf
import csv
import docx
import openpyxl
import pptx
import email
from aiogram import types
from bs4 import BeautifulSoup
from imaginepy import AsyncImagine, Style, Ratio

class botmedia:
    def __init__(self,HG_img2text):
        self.model = whisper.load_model('tiny')
        self.HG_img2text = HG_img2text
    async def generate_imagecaption(self, url, HG_TOKEN):
        headers = {"Authorization": f"Bearer {HG_TOKEN}"}
        retries = 0
        async with aiohttp.ClientSession() as session:
            async with session.get(url) as resp1:
                if resp1.status != 200:
                    return f"Error: failed to download image ({resp1.status})"
                async with session.post(self.HG_img2text, headers=headers, data=await resp1.read()) as resp2:
                    if resp2.status == 200:
                        response = await resp2.json()
                        return "This image looks like a " + response[0]['generated_text']
                    elif resp2.status >= 500 or 'loading' in (await resp2.text()).lower():
                        retries += 1
                        if retries <= 3:
                            await asyncio.sleep(3)
                            return await self.generate_imagecaption(url,HG_TOKEN)
                        else:
                            return f"Server error: {await resp2.text()}"
                    else:
                        return f"Error: {await resp2.text()}"

    async def generate_image(self,image_prompt, style_value, ratio_value, negative):
        
        filename = "image.png"
        style_enum = Style[style_value]
        ratio_enum = Ratio[ratio_value]
        imagine = AsyncImagine(style_enum)
        img_data = await imagine.sdprem(
            prompt=image_prompt,
            style=style_enum,
            ratio=ratio_enum,
            priority="1",
            high_res_results="1",
            steps="70",
            negative=negative
        )
        if img_data is None:
            print("An error occurred while generating the image.")
            return
        #img_data = await imagine.upscale(image=img_data)
        try:
            with open(filename, mode="wb") as img_file:
                img_file.write(img_data)
        except Exception as e:
            print(f"An error occurred while writing the image to file: {e}")
            return None
        await imagine.close()
        return filename
    
    async def transcribe_audio(self, audio_file_path):
        with open(audio_file_path, 'rb') as audio_file:
            content = audio_file.read()

        result = self.model.transcribe(audio_file_path)
        transcription = result["text"]
        return transcription
    
    
    async def download_file(self,message: types.Message):
        if message.audio is not None:
            file = message.audio
            file_extension = file.file_name.split(".")[-1] if file.file_name is not None else 'ogg'
        elif message.voice is not None:
            file = message.voice
            file_extension = 'ogg'
        elif message.document is not None:
            file = message.document
            file_extension = file.file_name.split(".")[-1] if file.file_name is not None else ''
        else:
            return None
        file_path = f'{file.file_id}.{file_extension}'
        file_dir = 'downloaded_files'
        os.makedirs(file_dir, exist_ok=True)
        full_file_path = os.path.join(file_dir, file_path)
        downloadfile = await file.download(destination_file=full_file_path)
        return full_file_path
    
    async def read_document(self,filename):
        valid_extensions = ['txt', 'rtf', 'md', 'html', 'xml', 'csv', 'json', 'js', 'css', 'py', 'java', 'c', 'cpp', 'php', 'rb', 'swift', 'sql', 'sh', 'bat', 'ps1', 'ini', 'cfg', 'conf', 'log', 'svg', 'epub', 'mobi', 'tex', 'docx', 'odt', 'xlsx', 'ods', 'pptx', 'odp', 'eml', 'htaccess', 'nginx.conf', 'pdf']
        extension = filename.split('.')[-1]
        if extension not in valid_extensions:
            return 'Invalid document file'
        if extension == 'pdf':
            with open(filename, 'rb') as f:
                pdf_reader = pypdf.PdfReader(f)
                num_pages = len(pdf_reader.pages)
                contents = ''
                for page_num in range(num_pages):
                    page_obj = pdf_reader.pages[page_num]
                    page_text = page_obj.extract_text()
                    contents += page_text
            return contents
        elif extension == 'docx':
            doc = docx.Document(filename)
            contents = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            return contents
        elif extension in ['xlsx', 'ods']:
            workbook = openpyxl.load_workbook(filename, read_only=True)
            sheet = workbook.active
            contents = ''
            for row in sheet.iter_rows(values_only=True):
                contents += '\t'.join([str(cell_value) for cell_value in row]) + '\n'
            return contents
        elif extension in ['pptx', 'odp']:
            presentation = pptx.Presentation(filename)
            contents = ''
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        contents += shape.text + '\n'
            return contents
        elif extension == 'eml':
            with open(filename, 'r') as f:
                msg = email.message_from_file(f)
                contents = ''
                for part in msg.walk():
                    if part.get_content_type() == 'text/plain':
                        contents += part.get_payload()
            return contents
        elif extension in ['html', 'xml']:
            with open(filename, 'r') as f:
                soup = BeautifulSoup(f, 'html.parser')
                contents = soup.get_text()
            return contents
        elif extension == 'csv':
            with open(filename, 'r') as f:
                reader = csv.reader(f)
                contents = ''
                for row in reader:
                    contents += '\t'.join(row) + '\n'
            return contents

        else:
            with open(filename, 'r') as f:
                contents = f.read()
                return contents